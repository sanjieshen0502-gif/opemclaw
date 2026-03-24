import os
import sys
from pathlib import Path

# 获取桌面路径
desktop = Path.home() / "Desktop"
print(f"桌面路径: {desktop}")

# 列出所有文件
for f in desktop.iterdir():
    if "美的" in str(f.name):
        print(f"找到文件: {f.name} - 大小: {f.stat().st_size}字节")

# 尝试读取PPTX文件
import subprocess
import json

# 尝试使用python-pptx库
try:
    from pptx import Presentation
    print("\n尝试使用python-pptx库读取...")
    
    for f in desktop.iterdir():
        if f.suffix == '.pptx' and "美的" in str(f.name):
            print(f"\n正在读取: {f.name}")
            prs = Presentation(str(f))
            print(f"幻灯片数量: {len(prs.slides)}")
            
            for i, slide in enumerate(prs.slides):
                print(f"\n=== 幻灯片 {i+1} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            print(f"文本: {text}")
            break
except ImportError as e:
    print(f"python-pptx未安装: {e}")
    print("\n尝试使用Office COM对象（仅限Windows）...")
    
    try:
        import win32com.client
        import pythoncom
        
        for f in desktop.iterdir():
            if f.suffix == '.pptx' and "美的" in str(f.name):
                print(f"\n正在使用COM对象读取: {f.name}")
                pythoncom.CoInitialize()
                app = win32com.client.Dispatch("PowerPoint.Application")
                presentation = app.Presentations.Open(str(f))
                
                print(f"幻灯片数量: {presentation.Slides.Count}")
                for i in range(1, presentation.Slides.Count + 1):
                    slide = presentation.Slides(i)
                    print(f"\n=== 幻灯片 {i} ===")
                    # 提取文本
                    for shape in slide.Shapes:
                        if hasattr(shape, "TextFrame") and hasattr(shape.TextFrame, "TextRange"):
                            text = shape.TextFrame.TextRange.Text.strip()
                            if text:
                                print(f"文本: {text}")
                
                presentation.Close()
                app.Quit()
                pythoncom.CoUninitialize()
                break
                
    except ImportError as e2:
        print(f"win32com未安装: {e2}")
        print("\n备用方案: 使用pptx库（需要先安装）")