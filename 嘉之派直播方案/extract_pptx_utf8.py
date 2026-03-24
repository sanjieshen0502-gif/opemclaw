# -*- coding: utf-8 -*-
from pptx import Presentation
import sys
import os

# Find the PPTX file
inbound_dir = r"C:\Users\Komorebi\.openclaw\media\inbound"
output_dir = r"C:\Users\Komorebi\.openclaw\workspace\嘉之派直播方案"

# List files to find the pptx
pptx_file = None
for f in os.listdir(inbound_dir):
    if f.endswith('.pptx') and '美的' in f:
        pptx_file = os.path.join(inbound_dir, f)
        break

if not pptx_file:
    # Just get the most recent pptx
    pptx_files = [f for f in os.listdir(inbound_dir) if f.endswith('.pptx')]
    if pptx_files:
        pptx_files.sort(key=lambda x: os.path.getmtime(os.path.join(inbound_dir, x)), reverse=True)
        pptx_file = os.path.join(inbound_dir, pptx_files[0])

if not pptx_file or not os.path.exists(pptx_file):
    print(f"未找到PPTX文件")
    sys.exit(1)

print(f"读取文件: {pptx_file}")

prs = Presentation(pptx_file)
full_text = []
for i, slide in enumerate(prs.slides, 1):
    slide_text = [f'=== 第{i}页 ===']
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            slide_text.append(shape.text.strip())
    if len(slide_text) > 1:
        full_text.append('\n'.join(slide_text))

output = '\n\n'.join(full_text)

# Save to file
output_file = os.path.join(output_dir, '话术01_美的净水器脚本.md')
with open(output_file, 'w', encoding='utf-8') as f:
    f.write(f"# 话术01：美的净水器直播脚本\n")
    f.write(f"_来源：{os.path.basename(pptx_file)}\n")
    f.write(f"_提取时间：2026-03-24_\n\n")
    f.write("---\n\n")
    f.write(output)

print(f"已保存到: {output_file}")
