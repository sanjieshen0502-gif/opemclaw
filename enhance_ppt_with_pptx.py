#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
使用python-pptx增强PPT痛点话术
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# 文件路径
ppt_path = r"C:\Users\Komorebi\Desktop\美的\美的净水器逐字稿（互动版）修改版.pptx"
output_path = r"C:\Users\Komorebi\.openclaw\workspace\美的净水器逐字稿（增强版）.pptx"

# 要添加的痛点话术（精选）
enhanced_pain_points = [
    "家里有宝宝的姐妹注意！给宝宝冲奶粉的水，你真的放心吗？自来水里的重金属、抗生素，哪怕只有一丁点，对宝宝娇嫩的肾脏都是负担！",
    "老人年纪大了，免疫力下降。长期喝不干净的水，容易引发慢性病。很多老人为了省钱，烧自来水喝了一辈子，肾结石、关节痛都找上门！",
    "喜欢喝茶喝咖啡的家人，你们肯定懂！几百块一斤的茶叶，用自来水一泡，茶汤浑浊、香气全无！好茶必须配好水！",
    "看看你家的烧水壶、加湿器、蒸汽熨斗，是不是一层厚厚的水垢？水垢不仅费电，还缩短家电寿命！",
    "住老小区、高楼层的要特别注意！你们喝的是'二次供水'，水在楼顶水箱存了一夜，铁锈、泥沙、微生物翻倍！",
    "你以为桶装水就安全？很多小作坊直接灌自来水，桶身反复使用，细菌超标！而且一桶水喝一周，最后几天全是细菌繁殖的'重灾区'！"
]

def enhance_presentation():
    try:
        # 加载原PPT
        prs = Presentation(ppt_path)
        
        # 获取第一张幻灯片（索引0）
        slide1 = prs.slides[0]
        
        print(f"原PPT有 {len(prs.slides)} 张幻灯片")
        print(f"第一张幻灯片有 {len(slide1.shapes)} 个形状")
        
        # 分析第一张幻灯片的布局
        # 查找现有的文本框
        text_shapes = []
        for shape in slide1.shapes:
            if shape.has_text_frame:
                text_shapes.append(shape)
                print(f"文本框: '{shape.text[:50]}...'")
        
        # 在原幻灯片中添加新的文本框（痛点补充）
        # 在幻灯片底部添加
        
        # 计算位置（在现有内容下方）
        left = Inches(0.5)
        top = Inches(4.0)  # 在现有内容下方
        width = Inches(9.0)
        height = Inches(2.0)
        
        # 添加新文本框
        new_textbox = slide1.shapes.add_textbox(left, top, width, height)
        text_frame = new_textbox.text_frame
        
        # 设置文本框属性
        text_frame.word_wrap = True
        
        # 添加标题
        p = text_frame.paragraphs[0]
        p.text = "【更多健康痛点】"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 0, 0)  # 红色
        p.alignment = PP_ALIGN.LEFT
        
        # 添加痛点内容
        for i, point in enumerate(enhanced_pain_points[:3]):  # 只添加前3条，避免太多
            p = text_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0, 0, 0)  # 黑色
            p.space_after = Pt(6)
        
        # 保存增强后的PPT
        prs.save(output_path)
        
        print(f"\n增强完成！新文件已保存到: {output_path}")
        print(f"文件大小: {os.path.getsize(output_path) / 1024:.1f} KB")
        
        # 创建操作指南
        guide = f"""# PPT痛点增强操作指南

## 已完成的操作：
1. 已创建增强版PPT文件：{output_path}
2. 在第一张幻灯片底部添加了【更多健康痛点】文本框
3. 添加了3条精选痛点话术

## 如需进一步编辑：
1. 打开新PPT文件
2. 查看第一张幻灯片底部的红色标题文本框
3. 可调整文本框位置、大小或添加更多内容

## 增强的痛点话术内容：
{chr(10).join([f"{i+1}. {point}" for i, point in enumerate(enhanced_pain_points[:3])])}

## 备用痛点话术（可手动添加）：
{chr(10).join([f"{i+4}. {point}" for i, point in enumerate(enhanced_pain_points[3:])])}

## 注意事项：
- 原PPT文件未修改，已创建新文件
- 如需要更多痛点话术，可从'痛点话术补充.txt'文件中复制
- 建议根据直播时间选择3-5个最相关的痛点
"""
        
        guide_path = r"C:\Users\Komorebi\.openclaw\workspace\PPT痛点增强操作指南.txt"
        with open(guide_path, 'w', encoding='utf-8') as f:
            f.write(guide)
        
        print(f"操作指南已保存到: {guide_path}")
        
        return True
        
    except Exception as e:
        print(f"处理失败: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    if enhance_presentation():
        print("\n任务完成！")
    else:
        print("\n任务失败，请查看错误信息")